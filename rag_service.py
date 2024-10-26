import sys
import docx
import PyPDF2
import os
from dotenv import load_dotenv
from pptx import Presentation
from os import listdir
from os.path import isfile, join, isdir
from langchain_text_splitters import TokenTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams
from langchain_qdrant import Qdrant


# Função para listar arquivos de um diretório
def files_list(path):
    # Inicializando uma lista vazia
    file_list = []                 

    # Iterando sobre os arquivos do caminho especificado
    for f in listdir(path):
        full_path = join(path, f)
        # Se houver um arquivo, este deve ser adicionado à lista
        if isfile(full_path):
            file_list.append(full_path)
        # Se for uma pasta dentro do caminho especificado
        elif isdir(full_path):
            file_list += files_list(full_path)

    # A etapa acima permite que a função procure por arquivos em todas as possíveis pastas dentro do caminho especificado!
    return file_list

# O resultado da função acima será uma lista com os caminhos dos documentos encontrados.
# O passo seguinte será o tratamento dos documentos encontrados. Precisamos de funções
# capazes de ler e extrair informações destes.
# Inicialmente trabalharemos com documentos nos seguintes formatos: pdf, docx, txt e ppt.

# Para arquivos em word
def docx_read(file):
    doc = docx.Document(file)
    text = [part.text for part in doc.paragraphs]
    return '\n'.join(text)

# Para arquivos em pptx
def ppt_read(file):
    doc = Presentation(file)
    text = []
    for slide in doc.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return '\n'.join(text)

# Para extratir informações textuais dentro de pdfs e txts, podemos usar funções de python
# Agora criaremos a função que será responsável por todo o processamento dos documentos utilizados

def main_process(path):
    # Definição do modelo a ser utilizado para criar embeddings
    load_dotenv()
    # Configurações do modelo
    model_cfgs = {'device':'cpu'}
    # Configurações de codificação
    ecode_cfgs = {'normalize_embeddings':True}  # Para que todos sejam formatados na mesma escala

    # Inicializando a classe de embeddings
    pross = HuggingFaceEmbeddings(model_name = os.getenv("MODEL_NAME"),
                                  model_kwargs = model_cfgs,
                                  encode_kwargs = ecode_cfgs)
    
    # Agora temos o modelo responsável por processar os documento de interesse
    # Para este projeto usaremos o Qdrant, assim, precisamos inicializar o mesmo
    client = QdrantClient("http://localhost:6333")
    # Nome da coleção
    collection_name = "vdb-study"

    # Se a coleção já existir, a mesma deve ser excluida
    if client.collection_exists(collection_name):
        client.delete_collection(collection_name)

    # Uma nova coleção deve ser criada
    client.create_collection(collection_name,
                             vectors_config = VectorParams(size = 768, 
                                                           distance = Distance.DOT))
    
    # Agora temos total certeza de que a coleção foi criada, então passaremos a processar
    # os documentos e inserir as informações processadas no vector database.

    # Inicializando a instância Qdrant
    qdrant = Qdrant(client, collection_name, pross)

    print("Indexando documentos...")

    # Precisamos agora iterar sobre todos os arquivos identificados
    list_ = files_list(path)
    for file in list_:
        try:
            content = ""
            if file.endswith(".pdf"):
                
                reader = PyPDF2.PdfReader(file)
                for page in reader.pages:
                    content += " " + page.extract_text()

                print("Indexado:" + file)

            elif file.endswith(".txt"):
                with open(file,'r') as f:
                    content = f.read()

                print("Indexado:" + file)

            elif file.endswith(".docx"): 
                content = docx_read(file)
                print("Indexado:" + file)

            elif file.endswith(".pptx"): 
                content = ppt_read(file)
                print("Indexado:" + file)

            else:
                continue

            text_spliter = TokenTextSplitter(chunk_size = 500,
                                             chunk_overlap = 50)
            
            texts = text_spliter.split_text(content)

            # Usaremos os metadados para que possamos utilizar estes como
            # referência no output do LLM.
            metadata = [{"path":file} for _ in texts]

            qdrant.add_texts(texts, metadatas = metadata)
            
            
        except Exception as e:
            print(f"Error! - Details: {e}")

        print("Execução concluida!")


# Para verificar se o script está sendo executado dieramente
if __name__ == "__main__":
    args = sys.argv

    if len(args) > 1:
        main_process(args[1])
    else:
        print("Deve ser informado um caminho que contenha documentos!")